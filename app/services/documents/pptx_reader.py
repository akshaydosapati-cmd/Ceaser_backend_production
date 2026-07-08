from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.services.documents.schemas import ExtractedDocument


class PPTXReader:
    def read(self, path: Path) -> ExtractedDocument:
        presentation = Presentation(str(path))
        slides = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"Slide {index}\n" + "\n".join(texts))
        return ExtractedDocument(title=path.stem, pages=len(presentation.slides), content="\n\n".join(slides), metadata={"reader": "python-pptx"})
